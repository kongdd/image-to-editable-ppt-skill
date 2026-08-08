import sys
import tempfile
import unittest
import zipfile
from pathlib import Path

from PIL import Image
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
RUNTIME_DIR = ROOT / "skills/image-to-editable-ppt/cli/editppt/runtime"
sys.path.insert(0, str(RUNTIME_DIR))

from build_pptx_from_manifest import (  # noqa: E402
    content_box_for_manifest,
    emu,
    normalize_manifest,
    powerpoint_ooxml_issues,
    px_to_inches,
    render_preview,
    slide_size_type,
    text_box_xml,
    write_deck,
    write_pptx,
)
from prepare_deck_run import fit_content_box, slide_for_source  # noqa: E402


def scalable_test_font():
    candidates = (
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
    )
    return next((path for path in candidates if Path(path).exists()), None)


def preview_ink_center(manifest):
    with tempfile.TemporaryDirectory() as tmp:
        preview_path = Path(tmp) / "preview.png"
        render_preview(manifest, Path(tmp) / "manifest.json", preview_path)
        image = Image.open(preview_path).convert("RGB")

    dark_pixels = [
        (x, y)
        for y in range(image.height)
        for x in range(image.width)
        if max(image.getpixel((x, y))) < 128
    ]
    if not dark_pixels:
        return None
    xs = [point[0] for point in dark_pixels]
    ys = [point[1] for point in dark_pixels]
    return (min(xs) + max(xs)) / 2, (min(ys) + max(ys)) / 2


class SlideLayoutTest(unittest.TestCase):
    def test_export_uses_standard_powerpoint_ooxml_package(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            Image.new("RGB", (320, 180), "#2f6f9f").save(root / "asset.png")
            manifest = {
                "source": {"width_px": 1600, "height_px": 900},
                "slide": {"width": 13.333, "height": 7.5, "background": "#ffffff"},
                "content_box": {"left": 0, "top": 0, "width": 13.333, "height": 7.5},
                "images": [{"path": "asset.png", "box_px": [100, 100, 400, 220]}],
                "shapes": [
                    {"type": "roundRect", "box_px": [600, 100, 400, 220], "source_corner_radius_px": 18, "fill": "#dbeafe"},
                    {"type": "ellipse", "box_px": [1050, 100, 160, 160], "fill": "#dcfce7", "stroke": "none"},
                    {"type": "line", "points_px": [600, 350, 1000, 350], "stroke": "#334155", "dash": "dash"},
                    {"type": "polygon", "polygon_px": [[1100, 350], [1250, 450], [1050, 450]], "fill": "#fef3c7"},
                ],
                "text_boxes": [
                    {
                        "runs": [
                            {"text": "Microsoft ", "bold": True},
                            {"text": "PowerPoint", "color": "#1d4ed8"},
                        ],
                        "box_px": [100, 500, 900, 100],
                        "font_size": 28,
                    }
                ],
            }
            manifest_path = root / "manifest.json"
            manifest_path.write_text("{}", encoding="utf-8")
            output = root / "standard.pptx"

            write_pptx(manifest, output, manifest_path)

            self.assertEqual([], powerpoint_ooxml_issues(output))
            self.assertEqual(1, len(Presentation(output).slides))
            with zipfile.ZipFile(output) as archive:
                names = set(archive.namelist())
                self.assertIn("ppt/presProps.xml", names)
                self.assertIn("ppt/viewProps.xml", names)
                self.assertIn("ppt/tableStyles.xml", names)
                self.assertIn("Microsoft", archive.read("docProps/app.xml").decode("utf-8"))

    def test_svg_assets_are_embedded_without_raster_fallback(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            (root / "formula.svg").write_text(
                '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="80"><text x="5" y="50">E=mc²</text></svg>',
                encoding="utf-8",
            )
            manifest = {
                "source": {"width_px": 1600, "height_px": 900},
                "slide": {"width": 13.333, "height": 7.5},
                "content_box": {"left": 0, "top": 0, "width": 13.333, "height": 7.5},
                "images": [{"path": "formula.svg", "box_px": [100, 100, 600, 240]}],
            }
            output = root / "svg-native.pptx"

            write_pptx(manifest, output, root / "manifest.json")

            self.assertEqual([], powerpoint_ooxml_issues(output))
            with zipfile.ZipFile(output) as archive:
                media = [name for name in archive.namelist() if name.startswith("ppt/media/")]
                self.assertEqual(1, len(media))
                self.assertTrue(media[0].endswith(".svg"))
                self.assertEqual((root / "formula.svg").read_bytes(), archive.read(media[0]))
                slide_xml = archive.read("ppt/slides/slide1.xml").decode("utf-8")
                self.assertIn("asvg:svgBlip", slide_xml)
                self.assertNotIn(".png", "\n".join(archive.namelist()))

    def test_standard_powerpoint_export_preserves_speaker_notes(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            manifest = {
                "slide": {"width": 13.333, "height": 7.5},
                "text_boxes": [{"text": "Slide", "left": 1, "top": 1, "width": 2, "height": 1}],
            }
            output = root / "notes.pptx"

            write_deck(
                {"slide": {"width": 13.333, "height": 7.5}},
                [{"manifest": manifest, "manifest_path": root / "manifest.json"}],
                output,
                [{"page_index": 1, "text": "Original speaker note"}],
            )

            deck = Presentation(output)
            self.assertEqual("Original speaker note", deck.slides[0].notes_slide.notes_text_frame.text)
            self.assertEqual([], powerpoint_ooxml_issues(output))

    def test_non_wide_source_uses_source_pixel_size(self):
        slide = slide_for_source(1536, 1024)

        self.assertEqual("source", slide["size_mode"])
        self.assertAlmostEqual(16, slide["width"])
        self.assertAlmostEqual(10.6666667, slide["height"])

        box = fit_content_box(1536, 1024, slide)
        self.assertAlmostEqual(0, box["left"])
        self.assertAlmostEqual(0, box["top"])
        self.assertAlmostEqual(slide["width"], box["width"])
        self.assertAlmostEqual(slide["height"], box["height"])

    def test_near_wide_source_snaps_to_standard_wide_without_stretching(self):
        slide = slide_for_source(1918, 1080)

        self.assertEqual("wide", slide["size_mode"])
        self.assertAlmostEqual(13.333, slide["width"])
        self.assertAlmostEqual(7.5, slide["height"])

        box = fit_content_box(1918, 1080, slide)
        self.assertAlmostEqual(7.5, box["height"])
        self.assertGreater(box["left"], 0)
        self.assertLess(box["width"], slide["width"])

    def test_pixel_coordinates_map_to_content_box(self):
        manifest = {
            "source": {"width_px": 1536, "height_px": 1024},
            "slide": {"width": 13.333, "height": 7.5},
        }

        box = content_box_for_manifest(manifest)
        self.assertAlmostEqual(1.0415, box["left"], places=3)
        self.assertAlmostEqual(0, box["top"])
        self.assertAlmostEqual(11.25, box["width"])
        self.assertAlmostEqual(7.5, box["height"])

        position = px_to_inches(manifest, 0, 0, 1536, 1024)
        self.assertAlmostEqual(box["left"], position["left"])
        self.assertAlmostEqual(box["top"], position["top"])
        self.assertAlmostEqual(box["width"], position["width"])
        self.assertAlmostEqual(box["height"], position["height"])

    def test_non_wide_presentation_size_is_custom(self):
        self.assertEqual("custom", slide_size_type(emu(16), emu(10.6666667)))
        self.assertEqual("wide", slide_size_type(emu(13.333), emu(7.5)))

    def test_text_font_size_is_clamped_to_source_box(self):
        manifest = {
            "source": {"width_px": 1600, "height_px": 900},
            "slide": {"width": 13.333, "height": 7.5},
            "content_box": {"left": 0, "top": 0, "width": 13.333, "height": 7.5},
            "text_boxes": [
                {
                    "text": "Dense label",
                    "box_px": [100, 100, 260, 24],
                    "font_size": 24,
                }
            ],
        }

        normalized = normalize_manifest(manifest)
        text_box = normalized["text_boxes"][0]

        self.assertLess(text_box["font_size"], 24)
        self.assertEqual(24, text_box["_requested_font_size"])

    def test_text_font_size_fit_can_be_disabled(self):
        manifest = {
            "source": {"width_px": 1600, "height_px": 900},
            "slide": {"width": 13.333, "height": 7.5},
            "content_box": {"left": 0, "top": 0, "width": 13.333, "height": 7.5},
            "text_boxes": [
                {
                    "text": "Locked label",
                    "box_px": [100, 100, 260, 24],
                    "font_size": 24,
                    "fit_text": False,
                }
            ],
        }

        normalized = normalize_manifest(manifest)

        self.assertEqual(24, normalized["text_boxes"][0]["font_size"])

    def test_wrapped_text_fit_does_not_force_single_line_width(self):
        manifest = {
            "source": {"width_px": 1600, "height_px": 900},
            "slide": {"width": 13.333, "height": 7.5},
            "content_box": {"left": 0, "top": 0, "width": 13.333, "height": 7.5},
            "text_boxes": [
                {
                    "text": "Long body copy should wrap across multiple natural lines",
                    "box_px": [100, 100, 420, 120],
                    "font_size": 18,
                    "wrap": "square",
                }
            ],
        }

        normalized = normalize_manifest(manifest)

        self.assertGreater(normalized["text_boxes"][0]["font_size"], 10)
        self.assertLessEqual(normalized["text_boxes"][0]["font_size"], 18)

    def test_text_box_alignment_uses_drawingml_enum_values(self):
        xml = text_box_xml(
            2,
            {
                "text": "1",
                "left": 0,
                "top": 0,
                "width": 1,
                "height": 1,
                "align": "center",
                "valign": "middle",
            },
        )

        self.assertIn('algn="ctr"', xml)
        self.assertIn('anchor="ctr"', xml)
        self.assertNotIn('algn="center"', xml)
        self.assertNotIn('anchor="middle"', xml)

    def test_preview_centers_text_inside_its_box(self):
        manifest = {
            "source": {"width_px": 200, "height_px": 200},
            "slide": {"width": 2, "height": 2, "background": "#ffffff"},
            "content_box": {"left": 0, "top": 0, "width": 2, "height": 2},
            "preview_scale": 100,
            "text_boxes": [
                {
                    "text": "1",
                    "box_px": [50, 50, 100, 100],
                    "font_size": 48,
                    "fit_text": False,
                    "color": "#000000",
                    "align": "center",
                    "valign": "middle",
                }
            ],
        }

        ink_center = preview_ink_center(manifest)
        self.assertIsNotNone(ink_center)
        ink_center_x, ink_center_y = ink_center

        self.assertAlmostEqual(100, ink_center_x, delta=5)
        self.assertAlmostEqual(100, ink_center_y, delta=5)

    def test_preview_centers_mixed_size_runs_inside_their_box(self):
        preview_font = scalable_test_font()
        self.assertIsNotNone(preview_font)
        manifest = {
            "source": {"width_px": 200, "height_px": 200},
            "slide": {"width": 2, "height": 2, "background": "#ffffff"},
            "content_box": {"left": 0, "top": 0, "width": 2, "height": 2},
            "preview_scale": 100,
            "text_boxes": [
                {
                    "runs": [
                        {"text": "1", "font_size": 72},
                        {"text": "A", "font_size": 12},
                    ],
                    "box_px": [50, 50, 100, 100],
                    "font_size": 18,
                    "fit_text": False,
                    "preview_font": preview_font,
                    "color": "#000000",
                    "align": "center",
                    "valign": "middle",
                }
            ],
        }

        ink_center = preview_ink_center(manifest)
        self.assertIsNotNone(ink_center)
        ink_center_x, ink_center_y = ink_center

        self.assertAlmostEqual(100, ink_center_x, delta=8)
        self.assertAlmostEqual(100, ink_center_y, delta=8)

    def test_preview_rotates_centered_text_around_the_box_center(self):
        manifest = {
            "source": {"width_px": 200, "height_px": 200},
            "slide": {"width": 2, "height": 2, "background": "#ffffff"},
            "content_box": {"left": 0, "top": 0, "width": 2, "height": 2},
            "preview_scale": 100,
            "text_boxes": [
                {
                    "text": "1",
                    "box_px": [50, 50, 100, 100],
                    "font_size": 48,
                    "fit_text": False,
                    "color": "#000000",
                    "align": "center",
                    "valign": "middle",
                    "rotation": 45,
                }
            ],
        }

        ink_center = preview_ink_center(manifest)
        self.assertIsNotNone(ink_center)
        ink_center_x, ink_center_y = ink_center

        self.assertAlmostEqual(100, ink_center_x, delta=8)
        self.assertAlmostEqual(100, ink_center_y, delta=8)


if __name__ == "__main__":
    unittest.main()
